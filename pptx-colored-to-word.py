import collections
import collections.abc
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from docx import Document

# Open the PowerPoint file
pptx_file = "example.pptx"
pr = Presentation(pptx_file)

# Initialize a new Word document
doc = Document()

# Initialize a dictionary to store the red text and their corresponding slide pages
red_text_dict = {}

# Loop through all the slides in the PowerPoint file
for i, slide in enumerate(pr.slides):
    # Initialize a list to store the red text in this slide
    slide_red_text = []
    # Loop through all the shapes in each slide
    for shape in slide.shapes:
        # Check if the shape is a text box and has red font color
        if shape.has_text_frame and shape.text_frame.paragraphs:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font_color = run.font.color
                    if font_color.type == MSO_COLOR_TYPE.RGB:
                        if font_color.rgb == (255, 0, 0):
                            # Extract the text with red font color
                            red_text = run.text
                            # Add the red text to the slide_red_text list
                            slide_red_text.append(red_text)
    # Add the slide_red_text to the red_text_dict with its slide page
    if slide_red_text:
        red_text_dict[i+1] = slide_red_text

# Format the red text with its slide page(s) and add to the Word document
for page, sentences in red_text_dict.items():
    # Add the page number to the Word document
    doc.add_paragraph(f"Page {page}")
    # Add each red sentence to the Word document
    for sentence in sentences:
        doc.add_paragraph(sentence)

    doc.add_paragraph(f"\n")

# Save the Word document
doc.save('highlight_text.docx')